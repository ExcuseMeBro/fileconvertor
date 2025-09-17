#!/usr/bin/env python3
"""
File Converter Script - Converts various file formats to PDF
Supports: .docx, .doc, .xlsx, .xls, .pptx, .ppt, .png
"""

import os
import shutil
from pathlib import Path
import subprocess
import sys

def install_requirements():
    """Install required packages"""
    packages = [
        'docx2pdf',
        'python-docx',
        'openpyxl',
        'python-pptx',
        'Pillow',
        'reportlab',
        'xlsxwriter',
        'pypandoc'
    ]
    
    for package in packages:
        try:
            subprocess.check_call([sys.executable, '-m', 'pip', 'install', package])
            print(f"✓ Installed {package}")
        except subprocess.CalledProcessError:
            print(f"⚠ Failed to install {package}")

def create_directory_structure(source_dir, target_dir):
    """Create the same directory structure in target directory"""
    for root, dirs, files in os.walk(source_dir):
        # Calculate relative path from source
        rel_path = os.path.relpath(root, source_dir)
        
        # Create corresponding directory in target
        if rel_path != '.':
            target_path = os.path.join(target_dir, rel_path)
            os.makedirs(target_path, exist_ok=True)
            print(f"Created directory: {target_path}")

def convert_docx_to_pdf(input_file, output_file):
    """Convert DOCX/DOC to PDF using multiple methods"""
    try:
        # Method 1: Try docx2pdf
        from docx2pdf import convert
        convert(input_file, output_file)
        return True
    except ImportError:
        try:
            # Method 2: Try pypandoc
            import pypandoc
            pypandoc.convert_file(input_file, 'pdf', outputfile=output_file)
            return True
        except ImportError:
            try:
                # Method 3: Manual conversion using python-docx and reportlab
                from docx import Document
                from reportlab.lib.pagesizes import letter
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                from reportlab.lib.styles import getSampleStyleSheet
                
                doc_docx = Document(input_file)
                
                pdf_doc = SimpleDocTemplate(output_file, pagesize=letter)
                styles = getSampleStyleSheet()
                elements = []
                
                for paragraph in doc_docx.paragraphs:
                    if paragraph.text.strip():
                        elements.append(Paragraph(paragraph.text, styles['Normal']))
                        elements.append(Spacer(1, 6))
                
                pdf_doc.build(elements)
                return True
            except Exception as e:
                print(f"Error converting {input_file}: {e}")
                return False
    except Exception as e:
        print(f"Error converting {input_file}: {e}")
        return False

def convert_xlsx_to_pdf(input_file, output_file):
    """Convert XLSX/XLS to PDF"""
    try:
        import openpyxl
        from reportlab.lib.pagesizes import letter, A4
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
        from reportlab.lib import colors
        
        # Load workbook
        wb = openpyxl.load_workbook(input_file)
        
        # Create PDF
        doc = SimpleDocTemplate(output_file, pagesize=A4)
        elements = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            
            # Get data from sheet
            data = []
            for row in ws.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    data.append([str(cell) if cell is not None else '' for cell in row])
            
            if data:
                # Create table
                table = Table(data)
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 14),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                elements.append(table)
        
        doc.build(elements)
        return True
    except Exception as e:
        print(f"Error converting {input_file}: {e}")
        return False

def convert_pptx_to_pdf(input_file, output_file):
    """Convert PPTX/PPT to PDF"""
    try:
        from pptx import Presentation
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        
        prs = Presentation(input_file)
        
        doc = SimpleDocTemplate(output_file, pagesize=letter)
        styles = getSampleStyleSheet()
        elements = []
        
        for i, slide in enumerate(prs.slides):
            # Add slide number
            elements.append(Paragraph(f"Slide {i+1}", styles['Heading1']))
            elements.append(Spacer(1, 12))
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    elements.append(Paragraph(shape.text, styles['Normal']))
                    elements.append(Spacer(1, 6))
            
            elements.append(Spacer(1, 20))
        
        doc.build(elements)
        return True
    except Exception as e:
        print(f"Error converting {input_file}: {e}")
        return False

def convert_png_to_pdf(input_file, output_file):
    """Convert PNG to PDF"""
    try:
        from PIL import Image
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Image as RLImage
        
        # Open image
        img = Image.open(input_file)
        
        # Create PDF
        doc = SimpleDocTemplate(output_file, pagesize=letter)
        
        # Calculate size to fit page
        page_width, page_height = letter
        img_width, img_height = img.size
        
        # Scale image to fit page while maintaining aspect ratio
        scale = min(page_width / img_width, page_height / img_height) * 0.8
        new_width = img_width * scale
        new_height = img_height * scale
        
        # Create ReportLab image
        rl_img = RLImage(input_file, width=new_width, height=new_height)
        
        doc.build([rl_img])
        return True
    except Exception as e:
        print(f"Error converting {input_file}: {e}")
        return False

def convert_file(input_file, output_file):
    """Convert file based on its extension"""
    ext = Path(input_file).suffix.lower()
    
    if ext in ['.docx', '.doc']:
        return convert_docx_to_pdf(input_file, output_file)
    elif ext in ['.xlsx', '.xls']:
        return convert_xlsx_to_pdf(input_file, output_file)
    elif ext in ['.pptx', '.ppt']:
        return convert_pptx_to_pdf(input_file, output_file)
    elif ext == '.png':
        return convert_png_to_pdf(input_file, output_file)
    else:
        print(f"Unsupported file type: {ext}")
        return False

def main():
    source_dir = "/Users/bro/PROJECTS/fileconvertor/boltfiles"
    target_dir = "/Users/bro/PROJECTS/fileconvertor/converted"
    
    # Install requirements
    print("Installing required packages...")
    install_requirements()
    
    # Create target directory
    os.makedirs(target_dir, exist_ok=True)
    
    # Create directory structure
    print("Creating directory structure...")
    create_directory_structure(source_dir, target_dir)
    
    # Supported extensions
    supported_extensions = {'.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.png'}
    
    # Convert files
    converted_count = 0
    failed_count = 0
    
    print("\nStarting file conversion...")
    
    for root, dirs, files in os.walk(source_dir):
        for file in files:
            file_path = Path(root) / file
            ext = file_path.suffix.lower()
            
            if ext in supported_extensions:
                # Calculate relative path and create output path
                rel_path = os.path.relpath(root, source_dir)
                if rel_path == '.':
                    output_dir = target_dir
                else:
                    output_dir = os.path.join(target_dir, rel_path)
                
                # Create output filename with .pdf extension
                output_filename = file_path.stem + '.pdf'
                output_path = os.path.join(output_dir, output_filename)
                
                print(f"Converting: {file_path} -> {output_path}")
                
                if convert_file(str(file_path), output_path):
                    converted_count += 1
                    print(f"✓ Successfully converted: {file}")
                else:
                    failed_count += 1
                    print(f"✗ Failed to convert: {file}")
    
    print(f"\nConversion completed!")
    print(f"Successfully converted: {converted_count} files")
    print(f"Failed conversions: {failed_count} files")

if __name__ == "__main__":
    main()