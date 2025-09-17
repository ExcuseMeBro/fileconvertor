#!/usr/bin/env python3
"""
Improved File Converter Script - Converts various file formats to PDF
Supports: .docx, .doc, .xlsx, .xls, .pptx, .ppt, .png
With enhanced error handling and logging
"""

import os
import shutil
from pathlib import Path
import subprocess
import sys
import logging
from datetime import datetime

# Setup logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('conversion.log'),
        logging.StreamHandler()
    ]
)

def install_requirements():
    """Install required packages"""
    packages = [
        'docx2pdf',
        'python-docx',
        'openpyxl',
        'python-pptx',
        'Pillow',
        'reportlab',
        'xlsxwriter'
    ]
    
    for package in packages:
        try:
            result = subprocess.run([sys.executable, '-m', 'pip', 'install', package], 
                                  capture_output=True, text=True)
            if result.returncode == 0:
                logging.info(f"✓ Installed {package}")
            else:
                logging.warning(f"⚠ Failed to install {package}: {result.stderr}")
        except Exception as e:
            logging.warning(f"⚠ Failed to install {package}: {e}")

def create_directory_structure(source_dir, target_dir):
    """Create the same directory structure in target directory"""
    for root, dirs, files in os.walk(source_dir):
        rel_path = os.path.relpath(root, source_dir)
        
        if rel_path != '.':
            target_path = os.path.join(target_dir, rel_path)
            os.makedirs(target_path, exist_ok=True)
            logging.info(f"Created directory: {target_path}")

def convert_docx_to_pdf(input_file, output_file):
    """Convert DOCX/DOC to PDF using multiple fallback methods"""
    methods = [
        ("docx2pdf", convert_with_docx2pdf),
        ("python-docx + reportlab", convert_with_docx_reportlab),
        ("LibreOffice", convert_with_libreoffice)
    ]
    
    for method_name, method_func in methods:
        try:
            logging.info(f"Trying {method_name} for {input_file}")
            if method_func(input_file, output_file):
                logging.info(f"✓ Successfully converted {input_file} using {method_name}")
                return True
        except Exception as e:
            logging.warning(f"Failed with {method_name}: {e}")
            continue
    
    logging.error(f"✗ All methods failed for {input_file}")
    return False

def convert_with_docx2pdf(input_file, output_file):
    """Convert using docx2pdf library"""
    try:
        from docx2pdf import convert
        convert(input_file, output_file)
        return True
    except ImportError:
        raise Exception("docx2pdf not available")

def convert_with_docx_reportlab(input_file, output_file):
    """Convert using python-docx and reportlab"""
    try:
        from docx import Document
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import inch
        
        doc_docx = Document(input_file)
        pdf_doc = SimpleDocTemplate(output_file, pagesize=letter, 
                                  rightMargin=72, leftMargin=72,
                                  topMargin=72, bottomMargin=18)
        
        styles = getSampleStyleSheet()
        elements = []
        
        for paragraph in doc_docx.paragraphs:
            if paragraph.text.strip():
                # Handle different paragraph styles
                if paragraph.style.name.startswith('Heading'):
                    style = styles['Heading1']
                else:
                    style = styles['Normal']
                
                elements.append(Paragraph(paragraph.text, style))
                elements.append(Spacer(1, 6))
        
        # Handle tables
        for table in doc_docx.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text for cell in row.cells])
                if row_text.strip():
                    elements.append(Paragraph(row_text, styles['Normal']))
                    elements.append(Spacer(1, 3))
        
        pdf_doc.build(elements)
        return True
    except ImportError:
        raise Exception("python-docx or reportlab not available")

def convert_with_libreoffice(input_file, output_file):
    """Convert using LibreOffice command line"""
    try:
        output_dir = os.path.dirname(output_file)
        result = subprocess.run([
            'libreoffice', '--headless', '--convert-to', 'pdf',
            '--outdir', output_dir, input_file
        ], capture_output=True, text=True, timeout=60)
        
        if result.returncode == 0:
            # LibreOffice creates file with same name but .pdf extension
            base_name = os.path.splitext(os.path.basename(input_file))[0]
            generated_pdf = os.path.join(output_dir, f"{base_name}.pdf")
            
            if os.path.exists(generated_pdf) and generated_pdf != output_file:
                shutil.move(generated_pdf, output_file)
            
            return os.path.exists(output_file)
        return False
    except (subprocess.TimeoutExpired, FileNotFoundError):
        raise Exception("LibreOffice not available or timeout")

def convert_xlsx_to_pdf(input_file, output_file):
    """Convert XLSX/XLS to PDF with improved formatting"""
    try:
        import openpyxl
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, PageBreak
        from reportlab.lib import colors
        from reportlab.lib.units import inch
        
        wb = openpyxl.load_workbook(input_file, data_only=True)
        doc = SimpleDocTemplate(output_file, pagesize=landscape(A4))
        elements = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            
            # Get actual data range
            data = []
            for row in ws.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    row_data = []
                    for cell in row:
                        if cell is not None:
                            row_data.append(str(cell))
                        else:
                            row_data.append('')
                    data.append(row_data)
            
            if data:
                # Create table with better styling
                table = Table(data)
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ('FONTSIZE', (0, 1), (-1, -1), 8),
                ]))
                elements.append(table)
                
                if len(wb.sheetnames) > 1:
                    elements.append(PageBreak())
        
        doc.build(elements)
        return True
    except Exception as e:
        logging.error(f"Error converting Excel file {input_file}: {e}")
        return False

def convert_pptx_to_pdf(input_file, output_file):
    """Convert PPTX/PPT to PDF with better formatting"""
    try:
        from pptx import Presentation
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import inch
        
        prs = Presentation(input_file)
        doc = SimpleDocTemplate(output_file, pagesize=letter)
        styles = getSampleStyleSheet()
        elements = []
        
        for i, slide in enumerate(prs.slides):
            # Add slide header
            elements.append(Paragraph(f"Slide {i+1}", styles['Heading1']))
            elements.append(Spacer(1, 12))
            
            # Extract text from shapes
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            
            # Add content
            for content in slide_content:
                elements.append(Paragraph(content, styles['Normal']))
                elements.append(Spacer(1, 6))
            
            # Add page break except for last slide
            if i < len(prs.slides) - 1:
                elements.append(PageBreak())
        
        doc.build(elements)
        return True
    except Exception as e:
        logging.error(f"Error converting PowerPoint file {input_file}: {e}")
        return False

def convert_png_to_pdf(input_file, output_file):
    """Convert PNG to PDF with better quality"""
    try:
        from PIL import Image
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Image as RLImage
        
        # Open and process image
        with Image.open(input_file) as img:
            # Convert to RGB if necessary
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            # Create PDF
            doc = SimpleDocTemplate(output_file, pagesize=letter)
            page_width, page_height = letter
            
            # Calculate size to fit page while maintaining aspect ratio
            img_width, img_height = img.size
            scale = min((page_width * 0.8) / img_width, (page_height * 0.8) / img_height)
            new_width = img_width * scale
            new_height = img_height * scale
            
            # Create ReportLab image
            rl_img = RLImage(input_file, width=new_width, height=new_height)
            doc.build([rl_img])
        
        return True
    except Exception as e:
        logging.error(f"Error converting image file {input_file}: {e}")
        return False

def convert_file(input_file, output_file):
    """Convert file based on its extension"""
    ext = Path(input_file).suffix.lower()
    
    conversion_map = {
        '.docx': convert_docx_to_pdf,
        '.doc': convert_docx_to_pdf,
        '.xlsx': convert_xlsx_to_pdf,
        '.xls': convert_xlsx_to_pdf,
        '.pptx': convert_pptx_to_pdf,
        '.ppt': convert_pptx_to_pdf,
        '.png': convert_png_to_pdf
    }
    
    if ext in conversion_map:
        return conversion_map[ext](input_file, output_file)
    else:
        logging.warning(f"Unsupported file type: {ext}")
        return False

def main():
    source_dir = "/Users/bro/PROJECTS/fileconvertor/boltfiles"
    target_dir = "/Users/bro/PROJECTS/fileconvertor/converted"
    
    logging.info("Starting improved file conversion process...")
    
    # Install requirements
    logging.info("Installing required packages...")
    install_requirements()
    
    # Create target directory
    os.makedirs(target_dir, exist_ok=True)
    
    # Create directory structure
    logging.info("Creating directory structure...")
    create_directory_structure(source_dir, target_dir)
    
    # Supported extensions
    supported_extensions = {'.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.png'}
    
    # Collect all files first
    all_files = []
    for root, dirs, files in os.walk(source_dir):
        for file in files:
            file_path = Path(root) / file
            if file_path.suffix.lower() in supported_extensions:
                all_files.append(file_path)
    
    logging.info(f"Found {len(all_files)} files to convert")
    
    # Convert files
    converted_count = 0
    failed_count = 0
    failed_files = []
    
    for i, file_path in enumerate(all_files, 1):
        logging.info(f"Processing file {i}/{len(all_files)}: {file_path.name}")
        
        # Calculate relative path and create output path
        rel_path = os.path.relpath(file_path.parent, source_dir)
        if rel_path == '.':
            output_dir = target_dir
        else:
            output_dir = os.path.join(target_dir, rel_path)
        
        # Create output filename with .pdf extension
        output_filename = file_path.stem + '.pdf'
        output_path = os.path.join(output_dir, output_filename)
        
        # Skip if already exists and is not empty
        if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
            logging.info(f"Skipping {file_path.name} - already converted")
            converted_count += 1
            continue
        
        if convert_file(str(file_path), output_path):
            converted_count += 1
            logging.info(f"✓ Successfully converted: {file_path.name}")
        else:
            failed_count += 1
            failed_files.append(str(file_path))
            logging.error(f"✗ Failed to convert: {file_path.name}")
    
    # Summary
    logging.info(f"\n{'='*50}")
    logging.info(f"CONVERSION SUMMARY")
    logging.info(f"{'='*50}")
    logging.info(f"Total files found: {len(all_files)}")
    logging.info(f"Successfully converted: {converted_count}")
    logging.info(f"Failed conversions: {failed_count}")
    
    if failed_files:
        logging.info(f"\nFailed files:")
        for failed_file in failed_files:
            logging.info(f"  - {failed_file}")
    
    logging.info(f"Conversion log saved to: conversion.log")

if __name__ == "__main__":
    main()